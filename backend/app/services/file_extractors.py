"""Extract plain text from supported file types and URLs (documents only — no image/audio/video ingest)."""

from __future__ import annotations

import re
from collections.abc import Callable
from pathlib import Path
from urllib.parse import urlparse

import httpx
import trafilatura
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader


def _clean(text: str) -> str:
    text = re.sub(r"\s+", " ", text)
    return text.strip()


# Not ingested into RAG (use PDF, Office, text, URLs instead).
INGEST_BLOCKED_SUFFIXES: frozenset[str] = frozenset(
    {
        ".png",
        ".jpg",
        ".jpeg",
        ".gif",
        ".webp",
        ".bmp",
        ".tif",
        ".tiff",
        ".heic",
        ".heif",
        ".mp3",
        ".wav",
        ".m4a",
        ".flac",
        ".ogg",
        ".opus",
        ".aac",
        ".mp4",
        ".mov",
        ".webm",
        ".mkv",
        ".avi",
        ".mpeg",
        ".mpg",
    }
)


def ingest_blocked_reason(suffix: str) -> str | None:
    s = suffix.lower()
    if not s.startswith("."):
        s = f".{s}" if s else ""
    if s in INGEST_BLOCKED_SUFFIXES:
        return (
            "Image, audio, and video files are not ingested. "
            "Use PDF, Word, Excel, PowerPoint, plain text, CSV, or a URL."
        )
    return None


def extract_from_txt(path: Path) -> str:
    return _clean(path.read_text(encoding="utf-8", errors="replace"))


def extract_from_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    parts: list[str] = []
    for page in reader.pages:
        t = page.extract_text() or ""
        parts.append(t)
    return _clean("\n".join(parts))


def extract_from_docx(path: Path) -> str:
    doc = DocxDocument(str(path))
    return _clean("\n".join(p.text for p in doc.paragraphs))


def extract_from_xlsx(path: Path) -> str:
    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append(" | ".join(cells))
    return _clean("\n".join(parts))


def extract_from_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return _clean("\n".join(parts))


def extract_from_url(url: str, timeout: float = 25.0) -> str:
    parsed = urlparse(url)
    if parsed.scheme not in ("http", "https"):
        raise ValueError("Only http(s) URLs supported")
    downloaded = trafilatura.fetch_url(url, no_ssl=False)
    if downloaded:
        text = trafilatura.extract(downloaded, include_comments=False, include_tables=True)
        if text:
            return _clean(text)
    with httpx.Client(timeout=timeout, follow_redirects=True) as client:
        r = client.get(url, headers={"User-Agent": "EnterpriseCopilotBot/1.0"})
        r.raise_for_status()
    soup = BeautifulSoup(r.text, "lxml")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    return _clean(soup.get_text(separator="\n"))


EXTENSION_HANDLERS: dict[str, Callable[[Path], str]] = {
    ".txt": extract_from_txt,
    ".md": extract_from_txt,
    ".csv": extract_from_txt,
    ".pdf": extract_from_pdf,
    ".docx": extract_from_docx,
    ".xlsx": extract_from_xlsx,
    ".xls": extract_from_xlsx,
    ".pptx": extract_from_pptx,
    ".ppt": extract_from_pptx,
}


def extract_text(path: Path, mime: str | None = None) -> str:
    ext = path.suffix.lower()
    block = ingest_blocked_reason(ext)
    if block:
        raise ValueError(block)
    if ext in EXTENSION_HANDLERS:
        return EXTENSION_HANDLERS[ext](path)
    if mime:
        m = mime.lower()
        if "pdf" in m:
            return extract_from_pdf(path)
        if "word" in m or "officedocument.wordprocessingml" in m:
            return extract_from_docx(path)
        if "spreadsheet" in m or "excel" in m:
            return extract_from_xlsx(path)
        if "presentation" in m or "powerpoint" in m:
            return extract_from_pptx(path)
        if m.startswith("image/") or m.startswith("audio/") or m.startswith("video/"):
            raise ValueError(
                "Image, audio, and video MIME types are not ingested. "
                "Use document formats (PDF, Office, text) or a URL."
            )
        if m.startswith("text/"):
            return extract_from_txt(path)
    raise ValueError(f"Unsupported file type: {ext or mime}")
